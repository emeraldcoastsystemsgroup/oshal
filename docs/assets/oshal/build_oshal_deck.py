"""Generate the OSHAL overview deck as a native PowerPoint (.pptx).

Draws real, editable shapes — boxes, block arrows, connectors, and swimlanes —
on the OSHAL brand palette. Produces reference-architecture, process-flow, and
swimlane diagrams (not text), plus content/metrics slides. Run:
    python docs/assets/oshal/build_oshal_deck.py

Standards: PEP 8, type-hinted helpers. See STANDARDS.md.

Changelog:
    2026-06-08  Initial generator: 30-slide professional deck with vector diagrams,
                explicit dynamic-bot-registration and A2A slides.
"""
from __future__ import annotations

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn

# ---- brand palette (docs/assets/oshal/visual-identity.md) -------------------
NIGHT = RGBColor(0x0B, 0x10, 0x20)
PANEL = RGBColor(0x14, 0x1F, 0x3D)
PANEL2 = RGBColor(0x0E, 0x17, 0x30)
SECTION_BG = RGBColor(0x06, 0x0A, 0x18)
CYAN = RGBColor(0x7D, 0xD3, 0xFC)
GREEN = RGBColor(0x34, 0xD3, 0x99)
WHITE = RGBColor(0xF8, 0xFA, 0xFC)
MUTED = RGBColor(0x94, 0xA3, 0xB8)
SOFT = RGBColor(0xBA, 0xE6, 0xFD)
AMBER = RGBColor(0xF5, 0xC9, 0x6B)
FONT = "Segoe UI"

SW, SH = 13.333, 7.5
prs = Presentation()
prs.slide_width = Inches(SW)
prs.slide_height = Inches(SH)
BLANK = prs.slide_layouts[6]


# ---- text engine ------------------------------------------------------------
def render(tf, paras, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE, space=4):
    """Render paragraphs (list of run-lists; each run = (text,size,color,bold))."""
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    for m in ("margin_left", "margin_right"):
        setattr(tf, m, Inches(0.1))
    for m in ("margin_top", "margin_bottom"):
        setattr(tf, m, Inches(0.04))
    for i, runs in enumerate(paras):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(space)
        p.space_before = Pt(0)
        for (t, sz, c, b) in runs:
            r = p.add_run()
            r.text = t
            r.font.size = Pt(sz)
            r.font.color.rgb = c
            r.font.bold = b
            r.font.name = FONT


def textbox(s, x, y, w, h, paras, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    render(tb.text_frame, paras, align, anchor)
    return tb


def box(s, x, y, w, h, paras, fill=PANEL, border=CYAN, rnd=0.09,
        align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, line_w=1.25):
    shp = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                             Inches(x), Inches(y), Inches(w), Inches(h))
    try:
        shp.adjustments[0] = rnd
    except Exception:
        pass
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if border is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = border
        shp.line.width = Pt(line_w)
    shp.shadow.inherit = False
    if isinstance(paras, str):
        paras = [[(paras, 14, WHITE, True)]]
    render(shp.text_frame, paras, align, anchor)
    return shp


def band(s, x, y, w, h, paras, fill=PANEL, border=None, align=PP_ALIGN.LEFT):
    return box(s, x, y, w, h, paras, fill=fill, border=border, rnd=0.04,
               align=align, anchor=MSO_ANCHOR.MIDDLE)


def arrow_r(s, x, y, w, h=0.34, color=GREEN):
    a = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(x), Inches(y), Inches(w), Inches(h))
    a.fill.solid(); a.fill.fore_color.rgb = color; a.line.fill.background(); a.shadow.inherit = False
    try:
        a.adjustments[0] = 0.55; a.adjustments[1] = 0.5
    except Exception:
        pass
    return a


def arrow_d(s, x, y, h, w=0.34, color=GREEN):
    a = s.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, Inches(x), Inches(y), Inches(w), Inches(h))
    a.fill.solid(); a.fill.fore_color.rgb = color; a.line.fill.background(); a.shadow.inherit = False
    try:
        a.adjustments[0] = 0.5; a.adjustments[1] = 0.55
    except Exception:
        pass
    return a


def connect(s, x1, y1, x2, y2, color=GREEN, w=1.75, dashed=False, arrow=True, both=False):
    c = s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    c.line.color.rgb = color
    c.line.width = Pt(w)
    ln = c.line._get_or_add_ln()
    if dashed:
        ln.append(ln.makeelement(qn('a:prstDash'), {'val': 'dash'}))
    if both:
        ln.append(ln.makeelement(qn('a:headEnd'), {'type': 'triangle', 'w': 'med', 'len': 'med'}))
    if arrow:
        ln.append(ln.makeelement(qn('a:tailEnd'), {'type': 'triangle', 'w': 'med', 'len': 'med'}))
    return c


def label(s, x, y, w, text, size=11, color=MUTED, align=PP_ALIGN.CENTER, bold=False):
    return textbox(s, x, y, w, 0.3, [[(text, size, color, bold)]], align=align, anchor=MSO_ANCHOR.MIDDLE)


# ---- slide scaffolding ------------------------------------------------------
PAGENO = [0]


def newslide(section=False):
    s = prs.slides.add_slide(BLANK)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = SECTION_BG if section else NIGHT
    bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(SW), Inches(0.1))
    bar.fill.solid(); bar.fill.fore_color.rgb = GREEN; bar.line.fill.background(); bar.shadow.inherit = False
    PAGENO[0] += 1
    return s


def heading(s, text):
    textbox(s, 0.62, 0.32, 12.1, 0.7, [[(text, 28, CYAN, True)]], anchor=MSO_ANCHOR.TOP)
    rule = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.64), Inches(1.04), Inches(2.5), Inches(0.05))
    rule.fill.solid(); rule.fill.fore_color.rgb = GREEN; rule.line.fill.background(); rule.shadow.inherit = False


def footer(s, note=""):
    textbox(s, 0.62, 7.05, 9.0, 0.35, [[("OSHAL — Open Swarm Harness Agent LLM" + ("   ·   " + note if note else ""),
            10, MUTED, False)]], anchor=MSO_ANCHOR.MIDDLE)
    textbox(s, 12.4, 7.05, 0.7, 0.35, [[(str(PAGENO[0]), 10, MUTED, False)]], align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)


def bullets(s, items, x=0.7, y=1.45, w=12.0, size=18, gap=0.6, marker="▸"):
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(gap * len(items) + 0.4))
    tf = tb.text_frame; tf.word_wrap = True
    for i, it in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(10)
        if isinstance(it, tuple):
            head, sub = it
        else:
            head, sub = it, None
        r = p.add_run(); r.text = f"{marker}  "; r.font.size = Pt(size); r.font.color.rgb = GREEN; r.font.bold = True; r.font.name = FONT
        r = p.add_run(); r.text = head; r.font.size = Pt(size); r.font.color.rgb = WHITE; r.font.bold = True; r.font.name = FONT
        if sub:
            r = p.add_run(); r.text = "  —  " + sub; r.font.size = Pt(size - 3); r.font.color.rgb = MUTED; r.font.bold = False; r.font.name = FONT
    return tb


# ============================================================================
# SLIDES
# ============================================================================
def s_cover():
    import math
    s = newslide(section=True)
    cx, cy = SW / 2, 2.4
    R, nr = 0.6, 0.27
    nodes = []
    for k in range(7):
        ang = math.radians(-90 + k * (360 / 7))
        nodes.append((cx + 2.5 * math.cos(ang), cy + 1.45 * math.sin(ang), ang))
    # 1) connector lines from CORE EDGE to NODE EDGE (drawn first, under the shapes)
    for (nx, ny, ang) in nodes:
        sx, sy = cx + R * math.cos(ang), cy + R * math.sin(ang)
        ex, ey = nx - nr * math.cos(ang), ny - nr * math.sin(ang)
        connect(s, sx, sy, ex, ey, color=RGBColor(0x2B, 0x3C, 0x57), w=1.5, arrow=False)
    # 2) outer nodes
    for (nx, ny, ang) in nodes:
        n = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(nx - nr), Inches(ny - nr), Inches(2 * nr), Inches(2 * nr))
        n.fill.solid(); n.fill.fore_color.rgb = PANEL; n.line.color.rgb = CYAN; n.line.width = Pt(1.5); n.shadow.inherit = False
    # 3) core on top (covers any line ends)
    core = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(cx - R), Inches(cy - R), Inches(2 * R), Inches(2 * R))
    core.fill.solid(); core.fill.fore_color.rgb = GREEN; core.line.color.rgb = CYAN; core.line.width = Pt(2.5); core.shadow.inherit = False
    render(core.text_frame, [[("OSHAL", 13, NIGHT, True)]], PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE)
    textbox(s, 0.8, 4.4, 11.7, 2.6, [
        [("OSHAL", 50, CYAN, True)],
        [("Open Swarm Harness Agent LLM", 24, WHITE, True)],
        [("", 8, WHITE, False)],
        [("Any harness.   Any agent.   Any LLM.   One swarm.", 18, GREEN, True)],
        [("A live framework for agent-backed applications — apps, tools, agents, workflows and UI expand at runtime.", 13, MUTED, False)],
    ], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.TOP)


def s_story():
    s = newslide(); heading(s, "The story in 30 seconds")
    bullets(s, [
        ("You don't ask OSHAL a question — you hand it a ticket.", None),
        ("It spins up a swarm of specialist bots", "each a different vendor's runtime on a different model"),
        ("They collaborate over a mesh, share one workspace, return a finished deliverable", None),
        ("Every call is cost-tracked, observable, and governed by the platform", None),
        ("You add agents, tools, workflows, and whole apps while it's running", "no rebuild"),
    ], y=1.5, gap=0.82, size=19)
    box(s, 0.7, 6.0, 11.9, 0.8, [[("The gap in agent systems was never prompting.  The gap was runtime.", 18, GREEN, True)]],
        fill=PANEL2, border=GREEN)
    footer(s)


def s_problem():
    s = newslide(); heading(s, "Why most agent projects stall")
    textbox(s, 0.7, 1.3, 12, 0.5, [[("A demo is one prompt and one model. Production needs all of this at once:", 16, SOFT, False)]])
    left = ["Multiple specialized agents", "Per-agent tool permissions", "Runtime observability & cost", "Workflow lifecycle & gates", "Containerized execution"]
    right = ["App packaging & versioning", "Repeatable validation", "Vendor flexibility (not one model)", "A way to grow without redeploying", "Demo → production, not demo-only"]
    box(s, 0.7, 2.0, 5.9, 3.5, [[("Production demands", 16, CYAN, True)]] + [[("•  " + t, 16, WHITE, False)] for t in left],
        fill=PANEL2, border=CYAN, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP)
    box(s, 6.75, 2.0, 5.9, 3.5, [[("Where stacks break", 16, CYAN, True)]] + [[("•  " + t, 16, WHITE, False)] for t in right],
        fill=PANEL2, border=CYAN, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP)
    box(s, 0.7, 5.7, 11.9, 0.75, [[("Teams end up with a single chatbot-with-tools, or a brittle pile of scripts. Neither survives a real workload.", 16, MUTED, True)]],
        fill=NIGHT, border=None)
    footer(s)


def s_thesis():
    s = newslide(); heading(s, "The OSHAL thesis")
    textbox(s, 0.7, 1.25, 12, 0.5, [[("Treat agents as runtime infrastructure — not hardcoded helpers.", 18, GREEN, True)]])
    cards = [
        ("Any harness, any model — per bot", "One ticket can cross OpenAI → Anthropic → Google, coordinated as one swarm."),
        ("The persona IS the swarm", "Each bot's persona embeds its own quality gate — no separate reviewer by default."),
        ("The platform is live", "Apps, tools, agents and workflows register at runtime and appear in the swarm at once."),
    ]
    x = 0.7
    for (h, sub) in cards:
        box(s, x, 2.0, 3.86, 3.4, [[(h, 17, CYAN, True)], [("", 6, WHITE, False)], [(sub, 15, WHITE, False)]],
            fill=PANEL, border=CYAN, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP)
        x += 4.06
    footer(s)


def s_refarch():
    s = newslide(); heading(s, "Reference architecture")
    # swarm zone (dashed boundary)
    z = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(2.4), Inches(1.3), Inches(6.2), Inches(4.5))
    z.adjustments[0] = 0.03; z.fill.solid(); z.fill.fore_color.rgb = RGBColor(0x0A, 0x12, 0x28)
    z.line.color.rgb = GREEN; z.line.width = Pt(1.25); z.shadow.inherit = False
    zln = z.line._get_or_add_ln(); zln.append(zln.makeelement(qn('a:prstDash'), {'val': 'dash'}))
    textbox(s, 2.55, 1.36, 4.0, 0.3, [[("OSHAL swarm — one Docker image", 11, GREEN, True)]])
    # controller + chips
    box(s, 2.6, 1.62, 5.8, 1.36, [[("SWARM CONTROLLER", 13, CYAN, True)], [("orchestrates · never calls an LLM", 10, MUTED, False)]],
        fill=PANEL2, border=GREEN, anchor=MSO_ANCHOR.TOP)
    for i, t in enumerate(["Queue Manager", "Phase Routing", "Ticket Lifecycle", "Cost Rollup"]):
        box(s, 2.74 + i * 1.41, 2.2, 1.3, 0.64, [[(t, 9, WHITE, True)]], fill=PANEL, border=CYAN, rnd=0.12)
    arrow_d(s, 5.35, 2.98, 0.14)
    # mesh spine
    box(s, 2.6, 3.14, 5.8, 0.44, [[("Redis mesh   ·   XADD / XREADGROUP / XACK", 12, NIGHT, True)]], fill=GREEN, border=GREEN)
    # mixed-vendor bot nodes wired to the mesh
    for i, (h, tag) in enumerate([("Bot node", "Codex CLI · OpenAI"), ("Bot node", "Claude Code · Anthropic"), ("Bot node", "Gemini CLI · Google")]):
        bx = 2.62 + i * 1.96
        box(s, bx, 3.8, 1.8, 1.3, [[(h, 12, WHITE, True)], [("persona + quality gate", 9, MUTED, False)], [("", 4, MUTED, False)], [(tag, 10, CYAN, True)]],
            fill=PANEL, border=CYAN, anchor=MSO_ANCHOR.TOP)
        connect(s, bx + 0.9, 3.8, bx + 0.9, 3.58, color=GREEN, w=1.6, both=True)
    textbox(s, 2.6, 5.2, 5.8, 0.32, [[("Each bot picks its own harness · provider · model — mixed vendors, one swarm", 10, SOFT, False)]], align=PP_ALIGN.CENTER)
    # operator (entry)
    box(s, 0.4, 1.95, 1.75, 0.9, [[("Operator", 13, WHITE, True)], [("submits a ticket", 9.5, MUTED, False)]], fill=PANEL2, border=CYAN, anchor=MSO_ANCHOR.TOP)
    arrow_r(s, 2.16, 2.25, 0.42)
    # postgres
    box(s, 8.9, 1.55, 2.05, 0.78, [[("Postgres", 12, WHITE, True)], [("tickets · cost · work items", 9, MUTED, False)]], fill=PANEL, border=CYAN, anchor=MSO_ANCHOR.TOP)
    connect(s, 8.4, 1.94, 8.9, 1.94, color=CYAN, w=1.5, both=True)
    # deliverable (exit)
    box(s, 11.05, 2.62, 2.1, 0.98, [[("Deliverable", 13, CYAN, True)], [("artifacts · reports", 9.5, MUTED, False)], [("tested code", 9.5, MUTED, False)]], fill=PANEL2, border=CYAN, anchor=MSO_ANCHOR.TOP)
    connect(s, 8.6, 3.2, 11.05, 3.1, color=GREEN, w=2.0)
    label(s, 9.1, 2.78, 1.9, "results out", 9.5, MUTED, PP_ALIGN.LEFT)
    # chromadb
    box(s, 8.9, 3.85, 2.05, 0.78, [[("ChromaDB", 12, WHITE, True)], [("RAG / swarm memory", 9, MUTED, False)]], fill=PANEL, border=CYAN, anchor=MSO_ANCHOR.TOP)
    connect(s, 8.6, 4.24, 8.9, 4.24, color=CYAN, w=1.5, both=True)
    # remote / edge agent via A2A (below the zone)
    box(s, 2.6, 6.05, 2.9, 0.7, [[("Remote / edge agent", 12, WHITE, True)], [("via A2A · Headscale — local MCP on a PC/Mac", 9, MUTED, False)]], fill=PANEL, border=SOFT, anchor=MSO_ANCHOR.TOP)
    connect(s, 3.95, 5.8, 3.95, 6.05, color=SOFT, w=1.6, both=True, dashed=True)
    # harnesses / providers band
    box(s, 5.6, 6.05, 7.4, 0.7, [[("Harnesses: Cline · Codex CLI · Claude Code · Gemini CLI", 11, SOFT, True)],
        [("→  40+ providers (OpenAI · Anthropic · Google · Bedrock · Azure · self-hosted Ollama / LM Studio)", 9.5, MUTED, False)]],
        fill=PANEL2, border=GREEN, rnd=0.06, align=PP_ALIGN.LEFT)
    footer(s, "controller orchestrates · bot nodes execute · any harness, any provider")


def s_layers():
    s = newslide(); heading(s, "Four layers, one Docker image")
    rows = [
        ("L4 · UI", "Cockpit dashboard + code-server workspace", CYAN),
        ("L3 · Swarm controller", "queue manager · routing · mesh · ticketing — never calls an LLM", GREEN),
        ("L2 · Bot nodes", "own ALL LLM execution · heartbeats · per-call cost telemetry", GREEN),
        ("L1 · Harnesses", "Cline · Codex CLI · Claude Code · Gemini CLI (+ noop)", CYAN),
        ("L0 · Providers", "40+ — Anthropic, OpenAI, Google, Bedrock, Groq, Ollama…", CYAN),
    ]
    y = 1.45
    for (h, sub, c) in rows:
        band(s, 1.6, y, 10.1, 0.92, [[(h + "    ", 16, c, True), (sub, 13, WHITE, False)]],
             fill=PANEL, border=c, align=PP_ALIGN.LEFT)
        y += 1.04
    box(s, 0.45, 1.45, 1.0, 5.0, [[("one\nimage".replace("\n", " "), 12, NIGHT, True)]], fill=GREEN, border=None, rnd=0.12)
    textbox(s, 0.3, 6.55, 12, 0.4, [[("bot-entrypoint.sh reads BOT_RUNTIME to pick the role — same artifact, any layer.", 13, MUTED, False)]])
    footer(s)


def s_mixmode():
    s = newslide(); heading(s, "The differentiator: mix-mode swarms")
    textbox(s, 0.7, 1.2, 12, 0.5, [[("Every bot independently picks its harness, provider, and model — in one ticket:", 16, SOFT, False)]])
    box(s, 0.7, 2.7, 1.9, 1.2, [[("TICKET", 16, NIGHT, True)]], fill=GREEN, border=None)
    rows = [("code-developer", "codex-cli · OpenAI · gpt-5.x-codex"),
            ("code-reviewer", "claude-code · Anthropic · Sonnet 4.6"),
            ("research-bot", "gemini-cli · Google · Gemini 2.5 Pro"),
            ("test-engineer", "cline · OpenAI-native")]
    for i, (b, cfg) in enumerate(rows):
        y = 1.95 + i * 0.92
        arrow_r(s, 2.75, y + 0.22, 0.7)
        box(s, 3.6, y, 8.9, 0.78, [[(b + "   ", 15, CYAN, True), (cfg, 13, WHITE, False)]],
            fill=PANEL, border=CYAN, align=PP_ALIGN.LEFT)
    box(s, 0.7, 6.0, 11.9, 0.7, [[("The dispatcher doesn't know or care which vendor is on the other side — the persona file is the contract.", 14, MUTED, True)]],
        fill=PANEL2, border=GREEN)
    footer(s)


def s_persona():
    s = newslide(); heading(s, "The persona IS the swarm")
    textbox(s, 0.7, 1.3, 12, 0.6, [[("By default, one workerBot per ticket type — its persona YAML carries the whole quality gate.", 17, WHITE, True)]])
    items = ["Mode A/B/C output classification (answer · artifact · escalation)",
             "Citation rules (RAG doc_id, web: provenance)",
             "A required artifact set (HANDOVER.md, RCA reports, scripts)",
             "A 5-section escalation packet when it can't self-gate"]
    bullets(s, items, y=2.2, gap=0.66, size=17)
    box(s, 0.7, 5.4, 11.9, 1.0, [[("No external reviewer bot by default — the persona reviews itself. Reviewer-grade output from one bot.", 16, GREEN, True)],
        [("(Manifests can still opt in to reviewerBot + maxRevisions when a worker genuinely can't self-gate.)", 12, MUTED, False)]],
        fill=PANEL2, border=GREEN)
    footer(s)


def flow_row(s, steps, y, x0=0.7, x1=12.6, bh=0.95, fill=PANEL, border=CYAN, numbered=True, start=1):
    n = len(steps)
    gap = 0.5
    bw = ((x1 - x0) - gap * (n - 1)) / n
    x = x0
    for i, st in enumerate(steps):
        paras = []
        if numbered:
            paras.append([(str(start + i), 13, GREEN, True)])
        if isinstance(st, tuple):
            paras.append([(st[0], 13, WHITE, True)])
            paras.append([(st[1], 11, MUTED, False)])
        else:
            paras.append([(st, 13, WHITE, True)])
        box(s, x, y, bw, bh, paras, fill=fill, border=border, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        if i < n - 1:
            arrow_r(s, x + bw + 0.06, y + bh / 2 - 0.17, gap - 0.12)
        x += bw + gap
    return bw


def s_lifecycle():
    s = newslide(); heading(s, "How work flows — the envelope lifecycle")
    flow_row(s, [("Poll", "approved tickets"), ("Decompose", "select agents per phase"), ("Dispatch", "XADD to mesh"), ("Pick up", "XREADGROUP")],
             y=1.7, bh=1.05)
    flow_row(s, [("Assemble", "persona + handovers"), ("Execute", "harness → provider"), ("ACK + record", "tokens + cost"), ("Advance", "next phase")],
             y=3.55, bh=1.05, start=5)
    box(s, 0.7, 5.25, 11.9, 0.95, [[("Channels: oshal:mesh:agent.{id}  (XADD / XREADGROUP / XACK)   ·   heartbeats: oshal:runtime-agent:{id} every 30s", 14, SOFT, True)],
        [("Cost lands in the chat_tasks table — the canonical $ source, not a log estimate.", 12, MUTED, False)]],
        fill=PANEL2, border=CYAN)
    footer(s)


def s_swimlane():
    s = newslide(); heading(s, "Swarm interaction — swimlane")
    lanes = ["Operator / UI", "Controller", "Redis mesh", "Bot node", "Workspace"]
    top, lane_h, lx, lw = 1.45, 1.02, 0.6, 2.1
    cont_x, cont_w = lx + lw + 0.05, SW - (lx + lw + 0.05) - 0.5
    for i, ln in enumerate(lanes):
        y = top + i * lane_h
        box(s, lx, y, lw, lane_h - 0.1, [[(ln, 12, NIGHT, True)]], fill=GREEN if i in (1, 3) else CYAN, border=None, rnd=0.06)
        band(s, cont_x, y, cont_w, lane_h - 0.1, [[("", 6, WHITE, False)]], fill=PANEL2 if i % 2 == 0 else NIGHT, border=None)
    # step boxes: (lane_index, col, label)
    cols = 6
    step_w = (cont_w - 0.4) / cols
    def cell(col):
        return cont_x + 0.2 + col * step_w

    def stepbox(lane, col, text, w=1):
        y = top + lane * lane_h + 0.12
        b = box(s, cell(col), y, step_w * w - 0.18, lane_h - 0.34, [[(text, 11, WHITE, True)]], fill=PANEL, border=CYAN, rnd=0.12)
        return (cell(col) + (step_w * w - 0.18) / 2, y + (lane_h - 0.34) / 2, cell(col), y, cell(col) + step_w * w - 0.18, y + lane_h - 0.34)
    s1 = stepbox(0, 0, "Submit ticket")
    s2 = stepbox(1, 1, "Queue + route")
    s3 = stepbox(2, 2, "Envelope XADD")
    s4 = stepbox(3, 3, "Persona + execute")
    s5 = stepbox(4, 4, "Write artifacts")
    s6 = stepbox(1, 5, "Advance / deliver")
    seq = [s1, s2, s3, s4, s5, s6]
    for a, b in zip(seq, seq[1:]):
        connect(s, a[4], a[1], b[2], b[1], color=GREEN, w=1.6)
    footer(s, "actors as lanes · steps flow left-to-right")


def s_mesh_a2a():
    s = newslide(); heading(s, "Two coordination paths: mesh vs A2A")
    # left: internal mesh
    box(s, 0.7, 1.4, 5.85, 0.6, [[("INSIDE the swarm — Redis mesh", 15, GREEN, True)]], fill=PANEL2, border=GREEN)
    box(s, 1.0, 2.2, 2.2, 0.8, [[("Controller", 13, WHITE, True)]], fill=PANEL, border=CYAN)
    box(s, 4.1, 2.2, 2.2, 0.8, [[("Bot node", 13, WHITE, True)]], fill=PANEL, border=CYAN)
    connect(s, 3.2, 2.45, 4.1, 2.45, color=GREEN); connect(s, 4.1, 2.75, 3.2, 2.75, color=CYAN)
    label(s, 1.0, 3.15, 5.3, "Redis Streams · XADD / XREADGROUP / XACK", 11, SOFT)
    box(s, 1.0, 3.65, 5.3, 1.95, [[("oshal:mesh:agent.{id}", 13, CYAN, True)],
        [("direct · alias · broadcast · capabilities", 12, WHITE, False)],
        [("consumer groups: swarm-execution / swarm-workers", 11, MUTED, False)],
        [("each bot consumes its own stream — no proxy bottleneck", 11, MUTED, False)]],
        fill=PANEL, border=CYAN, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP)
    # right: A2A
    box(s, 6.8, 1.4, 5.85, 0.6, [[("ACROSS machines — A2A", 15, AMBER, True)]], fill=PANEL2, border=AMBER)
    box(s, 7.1, 2.2, 2.2, 0.8, [[("Swarm", 13, WHITE, True)]], fill=PANEL, border=CYAN)
    box(s, 10.2, 2.2, 2.2, 0.8, [[("Remote / edge agent", 12, WHITE, True)]], fill=PANEL, border=AMBER)
    connect(s, 9.3, 2.6, 10.2, 2.6, color=AMBER)
    label(s, 7.1, 3.15, 5.3, "headscale-http · http · sse · stdio", 11, SOFT)
    box(s, 7.1, 3.65, 5.3, 1.95, [[("src/shared/types/a2a.ts", 13, AMBER, True)],
        [("intents: mcp.initialize · list-tools · call-tool · shutdown · status.sync", 12, WHITE, False)],
        [("remote-client MCP bridge over a Headscale overlay", 11, MUTED, False)],
        [("gateway is in-flight — see roadmap", 11, MUTED, False)]],
        fill=PANEL, border=AMBER, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP)
    box(s, 0.7, 5.85, 11.95, 0.62, [[("In-swarm bots talk over the mesh. External / remote agents join over A2A. Same swarm, two transports.", 14, GREEN, True)]],
        fill=NIGHT, border=None)
    footer(s)


def s_tooltiers():
    s = newslide(); heading(s, "Three tool tiers per bot")
    tiers = [("Tier 3 · LLM-embedded tools", "provider-native (e.g. built-in web search)", "emerging", MUTED),
             ("Tier 2 · Native harness tools", "Cline / Claude / Codex file, shell, browser", "built in", CYAN),
             ("Tier 1 · Shared OSHAL registry", "cross-harness · per-agent auth: auto / ask / off", "governed", GREEN)]
    y = 1.7
    for (h, sub, tag, c) in tiers:
        band(s, 1.5, y, 9.0, 1.2, [[(h, 17, c, True)], [(sub, 13, WHITE, False)]], fill=PANEL, border=c, align=PP_ALIGN.LEFT)
        box(s, 10.7, y + 0.2, 1.9, 0.8, [[(tag, 13, NIGHT, True)]], fill=c, border=None, rnd=0.2)
        y += 1.45
    textbox(s, 0.7, 6.2, 12, 0.5, [[("Tier 1 is governed: registry metadata → per-agent switch → runtime executor (built-in · CLI · API · MCP). Plus RAG (ChromaDB).", 13, MUTED, False)]])
    footer(s)


def s_build_pipeline():
    s = newslide(); heading(s, "Build pipeline — 7 phases, complexity-gated")
    steps = ["Intake", "Planning", "Specialist", "Execution", "Testing", "Review", "Delivery"]
    n = len(steps); x0, x1 = 0.7, 12.6; gap = 0.28
    bw = ((x1 - x0) - gap * (n - 1)) / n; x = x0
    for i, st in enumerate(steps):
        c = AMBER if st in ("Specialist", "Testing", "Review") else CYAN
        box(s, x, 2.4, bw, 1.2, [[("P%d" % (i + 1), 12, GREEN, True)], [(st, 13, WHITE, True)]], fill=PANEL, border=c)
        if i < n - 1:
            arrow_r(s, x + bw + 0.02, 2.9, gap - 0.06)
        x += bw + gap
    box(s, 0.7, 4.1, 11.9, 0.95, [[("Complexity gate gives each ticket only the phases it needs:", 14, SOFT, True)],
        [("low → intake·planning·execution·delivery (4)    medium → +testing·review    high → all 7 (+ optional Phase-8 architecture pre-round)", 13, WHITE, False)]],
        fill=PANEL2, border=AMBER, align=PP_ALIGN.LEFT)
    box(s, 0.7, 5.25, 11.9, 0.42, [[("Incident RCA pipeline — 3-phase, single persona-gated bot:", 14, SOFT, True)]],
        fill=NIGHT, border=None, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE)
    flow_row(s, ["Worker (investigate)", "Review", "Revision (if needed)"], y=5.8, x0=0.7, x1=9.5, bh=0.62, fill=PANEL, border=GREEN, numbered=False)
    footer(s)


def s_twoways():
    s = newslide(); heading(s, "Two ways to build on OSHAL")
    textbox(s, 0.7, 1.2, 12, 0.5, [[("Both produce the same thing — a registered swarm app.", 16, SOFT, False)]])
    box(s, 0.7, 1.95, 5.85, 4.2, [[("A · Build IN the framework", 18, CYAN, True)], [("", 6, WHITE, False)],
        [("You hand-author a swarm-app manifest:", 14, WHITE, True)],
        [("bots · tools · UI ribbon · routes · workflow", 13, MUTED, False)], [("", 6, WHITE, False)],
        [("You design the swarm.", 14, GREEN, True)], [("", 8, WHITE, False)],
        [("Example → Little Monsters", 15, SOFT, True)]],
        fill=PANEL, border=CYAN, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP)
    box(s, 6.75, 1.95, 5.85, 4.2, [[("B · Generate VIA packing", 18, GREEN, True)], [("", 6, WHITE, False)],
        [("An OSHAL agent interviews you and emits a", 14, WHITE, True)],
        [("complete single-purpose bot: persona + manifest + KB", 13, MUTED, False)], [("", 6, WHITE, False)],
        [("The framework builds the swarm.", 14, GREEN, True)], [("", 8, WHITE, False)],
        [("Examples → email-summarizer · Federal Capture · job-hunter", 15, SOFT, True)]],
        fill=PANEL, border=GREEN, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP)
    box(s, 0.7, 6.25, 11.9, 0.55, [[("Hand-craft for a multi-bot product. Pack when one business process should become an agent — fast.", 13, MUTED, True)]],
        fill=NIGHT, border=None)
    footer(s)


def s_littlemonsters():
    s = newslide(); heading(s, "Mode A — Little Monsters (built IN the framework)")
    textbox(s, 0.7, 1.2, 12, 0.5, [[("A voice-first ADHD study companion for K-12 — one manifest, a full product surface.", 16, SOFT, False)]])
    bullets(s, ["Six education bots (lecture-scribe, tutor, quiz-master, librarian, study-coach, writing-coach)",
                "Custom ribbon UI + per-class dynamic UI (one icon per active class)",
                "A new education workflow as its own ticket type",
                "RAG-grounded Socratic tutor (each class's textbook)",
                "Toggling the app activates/deactivates all six bots + ribbon — without stopping containers"],
            y=2.0, gap=0.66, size=16)
    box(s, 0.7, 6.0, 11.9, 0.6, [[("This is the “I want to design a multi-bot app” path.", 14, GREEN, True)]], fill=PANEL2, border=GREEN)
    footer(s)


def s_packing_flow():
    s = newslide(); heading(s, "Mode B — packing: a process becomes a bot")
    box(s, 0.7, 1.7, 2.3, 1.0, [[("Operator", 14, WHITE, True), ], [("~8-question interview", 11, MUTED, False)]], fill=PANEL, border=CYAN)
    arrow_r(s, 3.05, 2.03, 0.6)
    box(s, 3.75, 1.7, 2.4, 1.0, [[("codex-packer", 14, GREEN, True)], [("chat agent", 11, MUTED, False)]], fill=PANEL2, border=GREEN)
    arrow_r(s, 6.25, 2.03, 0.6)
    box(s, 6.95, 1.35, 5.7, 1.75, [[("emits 3 artifacts", 13, CYAN, True)],
        [("persona YAML  (the quality gate)", 12, WHITE, False)],
        [("swarm-apps/<slug>.yaml  (ticket type + UI)", 12, WHITE, False)],
        [("optional kb/ corpus", 12, WHITE, False)]],
        fill=PANEL, border=CYAN, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP)
    arrow_d(s, 9.6, 3.15, 0.5)
    box(s, 6.95, 3.75, 5.7, 0.8, [[("POST /api/swarm/apps/load", 13, GREEN, True)]], fill=PANEL2, border=GREEN)
    arrow_d(s, 9.6, 4.6, 0.5)
    box(s, 6.95, 5.15, 5.7, 0.8, [[("a new single-purpose bot — live in the cockpit", 13, WHITE, True)]], fill=GREEN, border=None)
    box(s, 0.7, 3.4, 5.7, 2.7, [[("The business process IS the agent", 15, CYAN, True)], [("", 6, WHITE, False)],
        [("Predefined inputs, baked-in rules, correct by construction. The same idea as one-bot-per-ticket-type, applied deliberately to a whole workflow.", 14, WHITE, False)]],
        fill=PANEL2, border=CYAN, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP)
    footer(s)


def s_packing_examples():
    s = newslide(); heading(s, "Mode B in the wild")
    cards = [("email-summarizer", "Literally emitted by codex-packer. Daily Gmail + Calendar digest for one Workspace. Read-only by default."),
             ("Federal Capture", "Packs an entire gov-contracting capture → proposal pipeline into ONE capture-specialist bot. RFP in → bid/no-bid + package."),
             ("job-hunter", "The same packing pattern shipped as a standalone product: framework-generated agents — OSHAL is the invisible how.")]
    x = 0.7
    for (h, sub) in cards:
        box(s, x, 1.7, 3.86, 3.7, [[(h, 16, GREEN, True)], [("", 8, WHITE, False)], [(sub, 14, WHITE, False)]],
            fill=PANEL, border=GREEN, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP)
        x += 4.06
    box(s, 0.7, 5.7, 11.9, 0.65, [[("Build-in gives you Little Monsters. Packing gives you a fleet of focused, self-contained bots.", 14, MUTED, True)]],
        fill=NIGHT, border=None)
    footer(s)


def s_dynamic_reg():
    s = newslide(); heading(s, "Dynamic bot registration (the live expansion loop)")
    textbox(s, 0.7, 1.2, 12, 0.45, [[("New capability enters a RUNNING swarm — no rebuild:", 16, SOFT, False)]])
    flow_row(s, [("Register tool", "/api/tools/runtime/register"), ("Create agent", "/api/swarm/agents"), ("Assign tool", "auth: auto/ask/off"), ("Generate persona", "YAML on disk")],
             y=1.85, bh=1.1)
    flow_row(s, [("Launch", "/api/agents/:id/launch"), ("Compose service", "oshal-bot · joins oshal net"), ("Heartbeat", "oshal:runtime-agent:{id}"), ("Mesh + registry", "subscribed & discoverable")],
             y=3.55, bh=1.1, start=5)
    box(s, 0.7, 5.3, 11.9, 1.0, [[("Static bots live in the registry (26) + a compose service, UUID-matched across compose · registry · Redis · Postgres.", 13, WHITE, False)],
        [("Benchmark: 18 tools + 18 agents + 18 containers created, healthy, mesh-subscribed in ~52s — zero residue on cleanup.", 14, GREEN, True)]],
        fill=PANEL2, border=GREEN, align=PP_ALIGN.LEFT)
    footer(s, "this loop is the line between a starter repo and a framework")


def s_appgallery():
    s = newslide(); heading(s, "The swim lanes — apps already on OSHAL")
    data = [("oshal-engineering", "build-in", "the build swarm + Workflow Studio UI"),
            ("issue-rca", "build-in", "domain RCA as a one-and-done ticket type"),
            ("little-monsters", "build-in", "6 bots + custom & dynamic ribbon UI"),
            ("capture-crm / federal-capture", "packing", "a full CRM pipeline packed into one bot"),
            ("email-summarizer", "packing", "codex-packer output, end to end"),
            ("codex-packer", "meta", "the bot that builds the bots")]
    # header
    cols = [(0.7, 4.0, "App"), (4.8, 2.2, "Mode"), (7.1, 5.5, "What it proves")]
    y = 1.5
    for (cx, cw, t) in cols:
        box(s, cx, y, cw, 0.5, [[(t, 13, NIGHT, True)]], fill=CYAN, border=None, rnd=0.04, align=PP_ALIGN.LEFT)
    y += 0.56
    for (app, mode, proof) in data:
        mc = GREEN if mode == "packing" else (AMBER if mode == "meta" else CYAN)
        box(s, 0.7, y, 4.0, 0.66, [[(app, 12, WHITE, True)]], fill=PANEL2, border=None, rnd=0.04, align=PP_ALIGN.LEFT)
        box(s, 4.8, y, 2.2, 0.66, [[(mode, 12, mc, True)]], fill=PANEL2, border=None, rnd=0.04, align=PP_ALIGN.LEFT)
        box(s, 7.1, y, 5.5, 0.66, [[(proof, 12, WHITE, False)]], fill=PANEL2, border=None, rnd=0.04, align=PP_ALIGN.LEFT)
        y += 0.72
    footer(s, "one framework contract, reused across very different products")


def s_numbers():
    s = newslide(); heading(s, "OSHAL by the numbers")
    metrics = [("4", "agent harnesses"), ("40+", "LLM providers"), ("26", "swarm bots (+19 full)"), ("68", "persona definitions"),
               ("31", "standard tools / 16 cats"), ("9", "example swarm apps"), ("~115", "services, full swarm"), ("~52s", "18 bots inserted")]
    x0, y0, w, h, gx, gy = 0.7, 1.6, 2.85, 1.7, 0.2, 0.35
    for i, (big, lab) in enumerate(metrics):
        col, row = i % 4, i // 4
        box(s, x0 + col * (w + gx), y0 + row * (h + gy), w, h,
            [[(big, 34, GREEN, True)], [(lab, 12, WHITE, False)]], fill=PANEL, border=CYAN)
    box(s, 0.7, 5.7, 11.9, 0.85, [[("$1.30 vs $4.05", 22, GREEN, True), ("   per real incident RCA  ", 15, WHITE, True), ("(−68%)", 18, CYAN, True)]],
        fill=PANEL2, border=GREEN)
    footer(s, "every number is from the repo")


def s_proof():
    s = newslide(); heading(s, "Proof — economics + live insertion")
    box(s, 0.7, 1.6, 5.85, 4.4, [[("The economics", 17, CYAN, True)], [("", 8, WHITE, False)],
        [("2-bot worker→reviewer pipeline", 15, WHITE, False)], [("$4.05", 30, MUTED, True)], [("", 6, WHITE, False)],
        [("OSHAL single persona-gated bot", 15, WHITE, False)], [("$1.30", 36, GREEN, True)], [("", 6, WHITE, False)],
        [("≈ 68% lower — reviewer-grade output from one bot", 14, CYAN, True)]],
        fill=PANEL, border=CYAN, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP)
    box(s, 6.75, 1.6, 5.85, 4.4, [[("The live-insertion benchmark", 17, GREEN, True)], [("", 8, WHITE, False)],
        [("18 runtime tools · 18 agents · 18 containers", 15, WHITE, False)],
        [("all healthy · heartbeating · registry-visible · mesh-subscribed", 13, MUTED, False)], [("", 8, WHITE, False)],
        [("~52s", 44, GREEN, True)], [("API → Redis → Docker · cleanup left zero residue", 13, MUTED, False)]],
        fill=PANEL, border=GREEN, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP)
    footer(s, "cost read from chat_tasks · benchmark repeatable")


def s_queues_sched():
    s = newslide(); heading(s, "Redis queues & schedulers")
    box(s, 0.7, 1.55, 5.85, 4.5, [[("Redis Streams mesh", 16, CYAN, True)], [("", 6, WHITE, False)],
        [("XADD / XREADGROUP / XACK · consumer groups", 13, WHITE, False)],
        [("at-least-once, horizontally scalable consumers", 12, MUTED, False)], [("", 8, WHITE, False)],
        [("Keys", 14, GREEN, True)],
        [("oshal:mesh:agent.{id}", 12, SOFT, False)],
        [("oshal:runtime-agent:{id}  (TTL 2h)", 12, SOFT, False)],
        [("oshal:scheduler", 12, SOFT, False)],
        [("oshal:subtask-registry:{parent}", 12, SOFT, False)]],
        fill=PANEL, border=CYAN, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP)
    box(s, 6.75, 1.55, 5.85, 4.5, [[("Cron scheduling + self-healing", 16, GREEN, True)], [("", 6, WHITE, False)],
        [("features/scheduling: cron-parser + runner + redis store", 13, WHITE, False)],
        [("surfaced as the agent-scheduler standard tool", 12, MUTED, False)], [("", 8, WHITE, False)],
        [("Autonomous timers", 14, CYAN, True)],
        [("heartbeats (30s) · hourly re-register", 12, SOFT, False)],
        [("stuck-agent watchdog", 12, SOFT, False)],
        [("agent-health-monitor", 12, SOFT, False)],
        [("stale mesh-channel cleanup", 12, SOFT, False)]],
        fill=PANEL, border=GREEN, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP)
    footer(s)


def s_deploy():
    s = newslide(); heading(s, "Deployment configs & scalability")
    tiers = [("Local / core", "11 services — cockpit/chat/API + Postgres/Redis/Chroma", CYAN, 3.4),
             ("Full local swarm", "~115 services — controller + all bot nodes", GREEN, 5.0),
             ("Kubernetes", "ops/deployment/kubernetes + any-bot-k8s · Headscale gateway", CYAN, 6.6)]
    for i, (h, sub, c, w) in enumerate(tiers):
        y = 1.7 + i * 1.25
        box(s, 0.9, y, w, 1.0, [[(h, 15, c, True)], [(sub, 12, WHITE, False)]], fill=PANEL, border=c, align=PP_ALIGN.LEFT)
        if i < 2:
            arrow_d(s, 1.4, y + 1.02, 0.22)
    box(s, 7.9, 1.7, 4.7, 4.05, [[("Scale levers", 16, GREEN, True)], [("", 6, WHITE, False)],
        [("per-bot containers — independent workers", 13, WHITE, False)],
        [("Stream consumer groups — add consumers", 13, WHITE, False)],
        [("one image, runtime switch", 13, WHITE, False)],
        [("K8s for elastic scale; watchdogs self-heal", 13, WHITE, False)], [("", 8, WHITE, False)],
        [("laptop → cluster · frontier → self-hosted", 13, CYAN, True)]],
        fill=PANEL2, border=GREEN, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP)
    footer(s)


def s_remote():
    s = newslide(); heading(s, "Remote client & edge agent")
    textbox(s, 0.7, 1.2, 12, 0.5, [[("A daemon on a PC/Mac runs a local MCP (stdio JSON-RPC), registers over Headscale — bidirectional:", 15, SOFT, False)]])
    # control plane
    box(s, 0.8, 2.2, 3.0, 1.2, [[("OSHAL control plane", 13, WHITE, True)], [("/api/remote-clients", 11, MUTED, False)]], fill=PANEL, border=CYAN)
    box(s, 5.3, 2.2, 2.7, 1.2, [[("Headscale", 13, AMBER, True)], [("private overlay", 11, MUTED, False)]], fill=PANEL2, border=AMBER)
    box(s, 9.5, 2.2, 3.0, 1.2, [[("Endpoint daemon", 13, WHITE, True)], [("+ local MCP server", 11, MUTED, False)]], fill=PANEL, border=CYAN)
    connect(s, 3.8, 2.6, 5.3, 2.6, color=GREEN); connect(s, 8.0, 2.6, 9.5, 2.6, color=GREEN)
    connect(s, 9.5, 3.05, 8.0, 3.05, color=CYAN); connect(s, 5.3, 3.05, 3.8, 3.05, color=CYAN)
    box(s, 0.8, 4.0, 5.8, 2.0, [[("Mode 1 — talks TO the swarm", 15, GREEN, True)], [("", 6, WHITE, False)],
        [("POST /api/remote-clients/:id/swarm/send", 12, SOFT, False)],
        [("(direct or broadcast) — your desktop participates in the swarm.", 13, WHITE, False)]],
        fill=PANEL, border=GREEN, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP)
    box(s, 6.8, 4.0, 5.8, 2.0, [[("Mode 2 — IS an agent of the swarm", 15, AMBER, True)], [("", 6, WHITE, False)],
        [("mesh/A2A envelope → bridge → remote-client queue → local MCP →", 12, WHITE, False)],
        [("result returns as a mesh reply. edge-agent = lightweight variant.", 12, WHITE, False)],
        [("intents: initialize · list-tools · call-tool · shutdown · status.sync", 11, MUTED, False)]],
        fill=PANEL, border=AMBER, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP)
    footer(s, "secured by private overlay + session/shared-secret; MCP stays on the endpoint")


def s_different():
    s = newslide(); heading(s, "Why it's different")
    rows = [("Multi-vendor runtimes in one swarm", "Yes", "No"),
            ("Runtime app / tool / agent registration", "Yes", "Rebuild"),
            ("Per-bot, per-call cost with vendor attribution", "Yes", "Rare"),
            ("Persona-as-quality-gate (no reviewer needed)", "Yes", "No"),
            ("Pack a business process into one bot", "Yes", "No")]
    box(s, 0.7, 1.5, 8.2, 0.55, [[("Capability", 13, NIGHT, True)]], fill=CYAN, border=None, rnd=0.04, align=PP_ALIGN.LEFT)
    box(s, 9.0, 1.5, 1.7, 0.55, [[("OSHAL", 13, NIGHT, True)]], fill=GREEN, border=None, rnd=0.04)
    box(s, 10.8, 1.5, 1.8, 0.55, [[("Typical", 13, NIGHT, True)]], fill=MUTED, border=None, rnd=0.04)
    y = 2.12
    for (cap, a, b) in rows:
        box(s, 0.7, y, 8.2, 0.74, [[(cap, 13, WHITE, False)]], fill=PANEL2, border=None, rnd=0.04, align=PP_ALIGN.LEFT)
        box(s, 9.0, y, 1.7, 0.74, [[(a, 14, GREEN, True)]], fill=PANEL, border=GREEN, rnd=0.06)
        box(s, 10.8, y, 1.8, 0.74, [[(b, 13, MUTED, True)]], fill=PANEL, border=None, rnd=0.06)
        y += 0.82
    footer(s)


def s_status():
    s = newslide(); heading(s, "Honest status")
    box(s, 0.7, 1.6, 5.85, 4.7, [[("Shipped", 17, GREEN, True)], [("", 6, WHITE, False)]] +
        [[("•  " + t, 14, WHITE, False)] for t in
         ["Mix-mode swarms, per-bot harness/model", "Swarm-app manifests + hot-load", "Harness packing (codex-packer)",
          "Dynamic tool/agent/bot insertion (Docker)", "Redis mesh, heartbeats, cost telemetry",
          "Build + incident pipelines", "Windows / Docker / K8s; local models"]],
        fill=PANEL, border=GREEN, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP)
    box(s, 6.75, 1.6, 5.85, 4.7, [[("In flight / planned", 17, AMBER, True)], [("", 6, WHITE, False)]] +
        [[("•  " + t, 14, WHITE, False)] for t in
         ["Workflow Studio compile-to-runtime", "Agentic Workflow Studio (describe → build)",
          "A2A gateway for external agents (prod proof)", "One-call create-and-start agent",
          "Haven persona-fronted home as default entry", "Repeatable local-LLM swarm profile"]],
        fill=PANEL, border=AMBER, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP)
    textbox(s, 0.7, 6.4, 12, 0.4, [[("Full detail: ROADMAP.md (today/target) and BACKLOG.md (done-when).", 12, MUTED, False)]])
    footer(s)


def s_summary():
    s = newslide(); heading(s, "Summary — the story, one slide")
    bullets(s, [("Problem", "agent demos die in production; the gap is runtime"),
                ("Thesis", "agents as runtime infrastructure; any harness/agent/LLM in one swarm; persona is the gate"),
                ("How", "controller orchestrates · bot nodes execute · mesh connects · cost tracked"),
                ("Build two ways", "hand-author in the framework (Little Monsters) or generate via packing (email-summarizer, Federal Capture, job-hunter)"),
                ("Live", "apps, tools, agents register at runtime — 18 bots in ~52s"),
                ("Proven & honest", "68% cheaper RCA; clear shipped-vs-roadmap")],
            y=1.5, gap=0.82, size=17)
    footer(s)


def s_close():
    s = newslide(section=True)
    textbox(s, 1.0, 2.4, 11.3, 1.6, [[("OSHAL makes agent swarms operational.", 38, CYAN, True)],
            [("Any harness.  Any agent.  Any LLM.  One swarm.", 20, GREEN, True)]], align=PP_ALIGN.CENTER)
    box(s, 1.8, 4.6, 9.7, 1.2, [[("Demo line", 12, GREEN, True)],
        [("“Register 18 tools, create 18 agents, launch 18 bot containers — discovered live in ~52 seconds, API to Redis to Docker.”", 15, WHITE, True)]],
        fill=PANEL, border=GREEN, align=PP_ALIGN.CENTER)
    textbox(s, 1.0, 6.1, 11.3, 0.5, [[("Ask: benchmark OSHAL on a real task workload — or pick a process and let codex-packer turn it into a live bot.", 13, MUTED, False)]], align=PP_ALIGN.CENTER)


for fn in [s_cover, s_story, s_problem, s_thesis, s_refarch, s_layers, s_mixmode, s_persona,
           s_lifecycle, s_swimlane, s_mesh_a2a, s_tooltiers, s_build_pipeline, s_twoways,
           s_littlemonsters, s_packing_flow, s_packing_examples, s_dynamic_reg, s_appgallery,
           s_numbers, s_proof, s_queues_sched, s_deploy, s_remote, s_different, s_status,
           s_summary, s_close]:
    fn()

OUT = "OSHAL-overview.pptx"
prs.save(OUT)
print(f"Saved {OUT} with {len(prs.slides._sldIdLst)} slides")
