"""OSHAL reference-architecture figure (v2) in Emerald Coast Systems Group brand colors.

A true reference architecture: ticket -> controller -> mixed-vendor bot swarm -> deliverable,
with the Redis mesh wired as the spine, Postgres/ChromaDB attached, and an A2A path to a
remote/edge agent. One-slide PPTX, exported to PNG via PowerPoint for the site + white paper.
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn

NAVY = RGBColor(0x0A, 0x1D, 0x30)
ZONE = RGBColor(0x0C, 0x22, 0x3A)
PANEL = RGBColor(0x0E, 0x2A, 0x47)
STRIP = RGBColor(0x0F, 0x2C, 0x49)
EMER = RGBColor(0x0E, 0x7C, 0x64)
BRIGHT = RGBColor(0x5F, 0xDC, 0xC0)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
MUT = RGBColor(0x9F, 0xB3, 0xC8)
SOFT = RGBColor(0xCD, 0xD9, 0xE6)
FONT = "Segoe UI"

SW, SH = 13.333, 7.5
prs = Presentation(); prs.slide_width = Inches(SW); prs.slide_height = Inches(SH)
s = prs.slides.add_slide(prs.slide_layouts[6])
s.background.fill.solid(); s.background.fill.fore_color.rgb = NAVY


def render(tf, paras, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE):
    tf.word_wrap = True; tf.vertical_anchor = anchor
    tf.margin_left = Inches(0.06); tf.margin_right = Inches(0.06)
    tf.margin_top = Inches(0.02); tf.margin_bottom = Inches(0.02)
    for i, runs in enumerate(paras):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align; p.space_after = Pt(1); p.space_before = Pt(0)
        for (t, sz, c, b) in runs:
            r = p.add_run(); r.text = t; r.font.size = Pt(sz); r.font.color.rgb = c; r.font.bold = b; r.font.name = FONT


def box(x, y, w, h, paras, fill=PANEL, border=BRIGHT, rnd=0.1, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, lw=1.25):
    sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    try: sh.adjustments[0] = rnd
    except Exception: pass
    sh.fill.solid(); sh.fill.fore_color.rgb = fill
    if border is None: sh.line.fill.background()
    else: sh.line.color.rgb = border; sh.line.width = Pt(lw)
    sh.shadow.inherit = False
    if isinstance(paras, str): paras = [[(paras, 13, WHITE, True)]]
    render(sh.text_frame, paras, align, anchor); return sh


def zone(x, y, w, h, labeltext):
    sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    sh.adjustments[0] = 0.03
    sh.fill.solid(); sh.fill.fore_color.rgb = ZONE
    sh.line.color.rgb = EMER; sh.line.width = Pt(1.25); sh.shadow.inherit = False
    ln = sh.line._get_or_add_ln(); ln.append(ln.makeelement(qn('a:prstDash'), {'val': 'dash'}))
    tb = s.shapes.add_textbox(Inches(x + 0.12), Inches(y + 0.06), Inches(3), Inches(0.3))
    render(tb.text_frame, [[(labeltext, 11, EMER, True)]], PP_ALIGN.LEFT, MSO_ANCHOR.TOP)
    return sh


def tbox(x, y, w, h, paras, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h)); render(tb.text_frame, paras, align, anchor); return tb


def arrow_r(x, y, w, h=0.3, color=BRIGHT):
    a = s.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(x), Inches(y), Inches(w), Inches(h))
    a.fill.solid(); a.fill.fore_color.rgb = color; a.line.fill.background(); a.shadow.inherit = False
    try: a.adjustments[0] = 0.55; a.adjustments[1] = 0.5
    except Exception: pass


def arrow_d(x, y, h, w=0.3, color=BRIGHT):
    a = s.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, Inches(x), Inches(y), Inches(w), Inches(h))
    a.fill.solid(); a.fill.fore_color.rgb = color; a.line.fill.background(); a.shadow.inherit = False
    try: a.adjustments[0] = 0.5; a.adjustments[1] = 0.55
    except Exception: pass


def link(x1, y1, x2, y2, color=BRIGHT, w=1.75, both=False, dashed=False):
    c = s.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    c.line.color.rgb = color; c.line.width = Pt(w)
    ln = c.line._get_or_add_ln()
    if dashed: ln.append(ln.makeelement(qn('a:prstDash'), {'val': 'dash'}))
    if both: ln.append(ln.makeelement(qn('a:headEnd'), {'type': 'triangle', 'w': 'med', 'len': 'med'}))
    ln.append(ln.makeelement(qn('a:tailEnd'), {'type': 'triangle', 'w': 'med', 'len': 'med'}))
    return c


# accent bar + title
bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(SW), Inches(0.1))
bar.fill.solid(); bar.fill.fore_color.rgb = EMER; bar.line.fill.background(); bar.shadow.inherit = False
tbox(0.55, 0.3, 12, 1.0, [[("OSHAL — reference architecture", 25, BRIGHT, True)],
     [("One ticket → a swarm of mixed-vendor agents → a delivered, cost-tracked result", 13, MUT, False)]])

# ---- the swarm zone ----
zone(2.4, 1.5, 6.2, 4.6, "OSHAL swarm — one Docker image")

# controller
box(2.6, 1.86, 5.8, 1.5, [[("SWARM CONTROLLER", 13, BRIGHT, True)], [("orchestrates · never calls an LLM", 10, MUT, False)]],
    fill=STRIP, border=EMER, anchor=MSO_ANCHOR.TOP)
for i, t in enumerate(["Queue\nManager", "Phase\nRouting", "Ticket\nLifecycle", "Cost\nRollup"]):
    box(2.74 + i * 1.41, 2.5, 1.3, 0.74, [[(t.replace("\n", " "), 9.5, WHITE, True)]], fill=PANEL, border=BRIGHT, rnd=0.12)
arrow_d(5.35, 3.36, 0.16)
# mesh spine
box(2.6, 3.52, 5.8, 0.46, [[("Redis mesh   ·   XADD / XREADGROUP / XACK", 12, WHITE, True)]], fill=EMER, border=EMER)
# bot nodes (mixed vendors)
bots = [("Bot node", "Codex CLI · OpenAI"), ("Bot node", "Claude Code · Anthropic"), ("Bot node", "Gemini CLI · Google")]
for i, (h, tag) in enumerate(bots):
    bx = 2.62 + i * 1.96
    box(bx, 4.18, 1.8, 1.4, [[(h, 12, WHITE, True)], [("persona + quality gate", 9, MUT, False)], [("", 4, MUT, False)], [(tag, 10.5, BRIGHT, True)]],
        fill=PANEL, border=BRIGHT, anchor=MSO_ANCHOR.TOP)
    link(bx + 0.9, 4.18, bx + 0.9, 3.98, color=BRIGHT, w=1.6, both=True)  # bot <-> mesh
tbox(2.6, 5.66, 5.8, 0.35, [[("Each bot picks its own harness · provider · model — mixed vendors, one swarm", 10.5, SOFT, False)]], align=PP_ALIGN.CENTER)

# ---- entry / exit / stores / remote ----
box(0.4, 2.2, 1.75, 0.95, [[("Operator", 13, WHITE, True)], [("submits a ticket", 9.5, MUT, False)]], fill=STRIP, border=BRIGHT, anchor=MSO_ANCHOR.TOP)
arrow_r(2.16, 2.5, 0.45)

box(8.95, 1.95, 2.05, 0.82, [[("Postgres", 12, WHITE, True)], [("tickets · cost · work items", 9, MUT, False)]], fill=PANEL, border=BRIGHT, anchor=MSO_ANCHOR.TOP)
link(8.6, 2.35, 8.95, 2.35, both=True, w=1.5)

box(8.95, 4.35, 2.05, 0.82, [[("ChromaDB", 12, WHITE, True)], [("RAG / swarm memory", 9, MUT, False)]], fill=PANEL, border=BRIGHT, anchor=MSO_ANCHOR.TOP)
link(8.6, 4.76, 8.95, 4.76, both=True, w=1.5)
tbox(8.62, 4.5, 0.4, 0.25, [[("RAG", 8.5, MUT, False)]])

box(11.15, 3.05, 2.0, 1.0, [[("Deliverable", 13, BRIGHT, True)], [("artifacts · reports", 9.5, MUT, False)], [("tested code", 9.5, MUT, False)]], fill=STRIP, border=BRIGHT, anchor=MSO_ANCHOR.TOP)
link(8.6, 3.62, 11.15, 3.55, w=2.0)
tbox(9.1, 3.18, 1.8, 0.25, [[("results out", 9.5, MUT, False)]])

box(2.6, 6.42, 2.7, 0.78, [[("Remote / edge agent", 12, WHITE, True)], [("runs a local MCP on a PC/Mac", 9, MUT, False)]], fill=PANEL, border=SOFT, anchor=MSO_ANCHOR.TOP)
link(3.95, 6.1, 3.95, 6.42, both=True, w=1.6, dashed=True, color=SOFT)
tbox(4.1, 6.12, 3.0, 0.25, [[("A2A · Headscale (http / sse / stdio)", 9.5, MUT, False)]])

# bottom band: harnesses / providers
box(5.5, 6.42, 7.6, 0.78, [[("Harnesses: Cline · Codex CLI · Claude Code · Gemini CLI", 11, SOFT, True)],
    [("→  40+ providers (OpenAI · Anthropic · Google · Bedrock · Azure · self-hosted Ollama / LM Studio)", 10, MUT, False)]],
    fill=STRIP, border=EMER, rnd=0.06, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE)

prs.save("_refarch_ecsg.pptx")
print("saved _refarch_ecsg.pptx")
