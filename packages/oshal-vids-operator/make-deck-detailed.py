"""Detailed OSHAL trade-recap deck from deck-data.json (DB + Alpaca extraction).
   More slides, the real analytical readout (results, YTD, why-the-triggers, sentiment,
   optimizations, book) — NOT time-boxed. Also emits a per-slide narration script for
   Vids "Convert Slides" (assign a voice + feed deck-narration.json slide by slide).
   Usage: python make-deck-detailed.py [deck-data.json]
   Out:   out/deck.pptx  +  out/deck-narration.json
"""
import json, os, sys
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE

HERE = os.path.dirname(os.path.abspath(__file__)); OUT = os.path.join(HERE, "out")
data = json.load(open(sys.argv[1] if len(sys.argv) > 1 else os.path.join(OUT, "deck-data.json"), encoding="utf-8"))

NAVY = RGBColor(0x07, 0x0D, 0x1F); CYAN = RGBColor(0x34, 0xE0, 0xFF)
GREEN = RGBColor(0x27, 0xE0, 0x8A); RED = RGBColor(0xFF, 0x5D, 0x6C)
WHITE = RGBColor(0xEA, 0xF2, 0xFF); MUTED = RGBColor(0x8A, 0xA0, 0xC8)
prs = Presentation(); prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
DATE = data.get("date", "")
narr = []  # one narration block per slide

def slide():
    s = prs.slides.add_slide(BLANK); s.background.fill.solid(); s.background.fill.fore_color.rgb = NAVY; return s
def box(s, x, y, w, h, text, size, color, bold=True, align=PP_ALIGN.LEFT, font="Arial", anchor=MSO_ANCHOR.TOP):
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h)); tf = tb.text_frame
    tf.word_wrap = True; tf.vertical_anchor = anchor
    for i, line in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align; r = p.add_run(); r.text = line
        f = r.font; f.size = Pt(size); f.bold = bold; f.color.rgb = color; f.name = font
    return tb
def rect(s, x, y, w, h, color):
    sh = s.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h)); sh.fill.solid(); sh.fill.fore_color.rgb = color; sh.line.fill.background(); return sh
def brand(s, kicker):
    rect(s, 0, 0, 13.333, 0.12, CYAN); box(s, 0.6, 0.3, 5, 0.55, "oshal", 26, CYAN)
    box(s, 2.3, 0.42, 6, 0.4, kicker, 13, MUTED)
    box(s, 9.0, 0.42, 3.8, 0.4, DATE, 13, MUTED, align=PP_ALIGN.RIGHT)
def money(n, dec=2):
    if n is None: return "—"
    return f"{'+' if n >= 0 else '-'}${abs(n):,.{dec}f}"

R = data.get("results", {}) or {}
Y = data.get("ytd") or {}
W = data.get("why", {}) or {}; WS = W.get("stats", {}) or {}; WT = W.get("themes", []) or []
SENT = data.get("sentiment", {}) or {}
TR = data.get("trades", []) or []
OPT = data.get("optimize", []) or []

# 1 — title
s = slide()
box(s, 0, 2.2, 13.333, 1.6, "oshal", 100, CYAN, align=PP_ALIGN.CENTER)
box(s, 0, 3.8, 13.333, 0.8, "DAILY TRADING REPORT", 28, WHITE, align=PP_ALIGN.CENTER)
box(s, 0, 4.6, 13.333, 0.6, DATE + "   ·   paper desk", 20, MUTED, align=PP_ALIGN.CENTER)
narr.append(f"Here is the Oh-shal daily trading report for {DATE}.")

# 2 — yesterday's results
s = slide(); brand(s, "YESTERDAY'S RESULTS")
pl = R.get("pl"); pct = R.get("pct")
box(s, 0.5, 1.3, 12, 1.8, money(pl), 90, GREEN if (pl or 0) >= 0 else RED)
box(s, 0.6, 3.1, 9, 0.6, f"Net day P/L   ·   {'+' if (pct or 0) >= 0 else ''}{pct}%", 22, MUTED)
box(s, 0.6, 4.4, 4, 0.4, "ACCOUNT EQUITY", 15, MUTED); box(s, 0.6, 4.8, 4, 0.7, f"${R.get('equity', 0):,.0f}", 36, WHITE)
box(s, 4.8, 4.4, 4, 0.4, "OPEN UNREALIZED", 15, MUTED); box(s, 4.8, 4.8, 4, 0.7, money(R.get("unrealized"), 0), 36, GREEN if (R.get("unrealized") or 0) >= 0 else RED)
box(s, 9.0, 4.4, 3.8, 0.4, "FILLS / POSITIONS", 15, MUTED); box(s, 9.0, 4.8, 3.8, 0.7, f"{R.get('fills', 0)}  /  {R.get('positions', 0)}", 36, CYAN)
narr.append(f"The desk finished {'up' if (pl or 0) >= 0 else 'down'} {abs(pl or 0):.0f} dollars, {abs(pct or 0)} percent, "
            f"with account equity at {R.get('equity', 0):,.0f} dollars across {R.get('positions', 0)} open positions and {R.get('fills', 0)} fills.")

# 3 — top winners (chart)
s = slide(); brand(s, "TOP OPEN WINNERS")
box(s, 0.6, 1.2, 11, 0.7, "Unrealized gains leading the book", 28, WHITE)
wins = R.get("winners", [])[:5]
if wins:
    cd = CategoryChartData(); cd.categories = [w[0] for w in wins]; cd.add_series("Unrealized %", tuple(float(w[1]) for w in wins))
    gf = s.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(0.6), Inches(2.4), Inches(12), Inches(4.4), cd)
    ch = gf.chart; ch.has_legend = False; ch.has_title = False; ch.plots[0].gap_width = 60
    ser = ch.plots[0].series[0]; ser.format.fill.solid(); ser.format.fill.fore_color.rgb = CYAN
    for ax in (ch.category_axis, ch.value_axis): ax.tick_labels.font.size = Pt(18); ax.tick_labels.font.color.rgb = WHITE
lead = ", ".join(f"{w[0]} up {w[1]} percent" for w in wins[:3])
narr.append(f"Leading the open book: {lead}." if lead else "The book is flat with no standout open winners.")

# 4 — SINCE INCEPTION performance (account has only traded a few days — NOT year-to-date)
s = slide(); brand(s, "SINCE INCEPTION")
if Y and Y.get("retPct") is not None:
    days = Y.get("days") or 0
    dlabel = f"{days} trading day{'s' if days != 1 else ''}"
    box(s, 0.5, 1.4, 12, 1.8, f"{'+' if Y['retPct'] >= 0 else ''}{Y['retPct']}%", 96, GREEN if Y["retPct"] >= 0 else RED)
    box(s, 0.6, 3.4, 12, 0.6, f"Total return  ·  {money(Y.get('retUsd'), 0)} over {dlabel}", 22, MUTED)
    box(s, 0.6, 4.6, 5, 0.4, "STARTING EQUITY", 15, MUTED); box(s, 0.6, 5.0, 5, 0.7, f"${Y.get('base', 0):,.0f}", 34, WHITE)
    box(s, 6.0, 4.6, 5, 0.4, "EQUITY AT REPORT", 15, MUTED); box(s, 6.0, 5.0, 5, 0.7, f"${Y.get('last', 0):,.0f}", 34, CYAN)
    narr.append(f"Since inception — just {dlabel} in — the desk is {'up' if Y['retPct'] >= 0 else 'down'} {abs(Y['retPct'])} percent, "
                f"about {abs(Y.get('retUsd') or 0):,.0f} dollars, from a {Y.get('base', 0):,.0f} dollar paper start.")
else:
    box(s, 0.6, 3, 12, 1, "Performance history unavailable.", 30, MUTED); narr.append("Performance history is not yet available for this account.")

# 5 — why the triggers fired
s = slide(); brand(s, "WHY THE TRIGGERS FIRED")
box(s, 0.6, 1.2, 12, 0.9, f"{WS.get('total', 0)} decisions  ·  {WS.get('buys', 0)} buys / {WS.get('sells', 0)} sells / {WS.get('holds', 0)} holds  ·  avg confidence {WS.get('avg_conf', 0)}", 22, CYAN)
y = 2.4
for t in WT[:6]:
    box(s, 0.8, y, 7, 0.5, t.get("theme", "").replace("_", " "), 24, WHITE)
    box(s, 8.2, y, 4.4, 0.5, f"{t.get('n', 0)} decisions   conf {t.get('avg_conf', 0)}", 20, MUTED)
    y += 0.66
themes = ", ".join(f"{t['theme'].replace('_',' ')} {t['n']}" for t in WT[:4])
narr.append(f"Every trade is justified to a signal. {WS.get('total', 0)} decisions today — {WS.get('buys', 0)} buys, {WS.get('sells', 0)} sells, {WS.get('holds', 0)} holds, "
            f"average confidence {int((WS.get('avg_conf') or 0)*100)} percent. The drivers: {themes}.")

# 6 — notable trades (the realized exits, by size)
s = slide(); brand(s, "NOTABLE TRADES")
realized = sorted([t for t in TR if t.get("pnl") is not None], key=lambda t: abs(t["pnl"]), reverse=True)[:7]
box(s, 0.6, 1.2, 12, 0.6, "Largest realized exits and why", 26, WHITE)
y = 2.3
box(s, 0.7, y, 2.2, 0.4, "SYMBOL", 14, MUTED); box(s, 2.9, y, 1.6, 0.4, "P&L", 14, MUTED); box(s, 4.7, y, 8, 0.4, "REASON", 14, MUTED); y += 0.5
for t in realized:
    box(s, 0.7, y, 2.2, 0.5, f"{t['symbol']} {t['side']}", 20, WHITE)
    box(s, 2.9, y, 1.8, 0.5, money(t["pnl"], 0), 20, GREEN if t["pnl"] >= 0 else RED)
    reason = (t.get("rationale") or "").split(" — ")[0].split(".")[0][:64]
    box(s, 4.7, y, 8.2, 0.5, reason, 16, MUTED)
    y += 0.58
top = realized[:2]
best = ", ".join(f"{t['symbol']} {money(t['pnl'],0)}" for t in sorted(realized, key=lambda t: t['pnl'], reverse=True)[:2])
worst = ", ".join(f"{t['symbol']} {money(t['pnl'],0)}" for t in sorted(realized, key=lambda t: t['pnl'])[:2])
narr.append(f"The day's biggest wins came from {best}; the biggest losses, {worst} — each one an automatic risk exit, a take profit, stop loss, or rotation.")

# 7 — market sentiment
s = slide(); brand(s, "MARKET SENTIMENT")
box(s, 0.6, 1.2, 12, 0.7, f"{SENT.get('totalSignals', 0)} market signals tracked — most active names", 24, CYAN)
tops = SENT.get("topSymbols", [])[:8]
box(s, 0.7, 2.1, 12, 0.8, "   ".join(f"{x['symbol']} ({x['mentions']})" for x in tops), 26, WHITE)
y = 3.3
for h in SENT.get("headlines", [])[:5]:
    box(s, 0.7, y, 12, 0.5, "•  " + (h.get("title") or "")[:96], 15, MUTED); y += 0.52
names = ", ".join(x["symbol"] for x in tops[:3])
narr.append(f"The signal stream tracked {SENT.get('totalSignals', 0)} market signals — the most active names were {names}, driven by the day's A-I, chip, and memory headlines.")

# 8 — optimizations (honest if empty)
s = slide(); brand(s, "OPTIMIZATIONS")
if OPT:
    box(s, 0.6, 1.2, 12, 0.6, "Parameter changes this session", 26, WHITE); y = 2.3
    for o in OPT[:7]:
        box(s, 0.7, y, 6, 0.5, o.get("param", ""), 20, WHITE)
        box(s, 6.8, y, 6, 0.5, f"{o.get('current_value')} → {o.get('proposed_value')}   ({o.get('status')})", 18, CYAN); y += 0.6
    narr.append(f"{len(OPT)} parameter optimizations were applied to tune expectancy and win rate.")
else:
    box(s, 0.6, 2.6, 12, 1.2, "No parameter changes recommended this session.", 34, MUTED)
    box(s, 0.6, 3.9, 12, 0.7, "The strategy held its current configuration.", 22, MUTED)
    narr.append("No parameter optimizations were recommended this session; the strategy held its current configuration.")

# 9 — book at close
s = slide(); brand(s, "BOOK AT CLOSE")
box(s, 0.6, 2.0, 6, 2.2, str(R.get("positions", 0)), 150, CYAN)
box(s, 0.6, 4.4, 6, 0.6, "OPEN POSITIONS", 20, MUTED)
box(s, 7.0, 2.3, 6, 0.5, "OPEN UNREALIZED", 16, MUTED)
box(s, 7.0, 2.7, 6, 1.0, money(R.get("unrealized"), 0), 50, GREEN if (R.get("unrealized") or 0) >= 0 else RED)
box(s, 7.0, 4.1, 6, 0.5, "LEADERS", 16, MUTED); box(s, 7.0, 4.5, 6, 0.8, (R.get("leaders") or "").replace(" · ", "   "), 32, WHITE)
narr.append(f"The book closed with {R.get('positions', 0)} open positions carrying about {abs(R.get('unrealized') or 0):,.0f} dollars in unrealized gains. That's the desk for {DATE}.")

deck_path = os.path.join(OUT, "deck.pptx"); prs.save(deck_path)
json.dump(narr, open(os.path.join(OUT, "deck-narration.json"), "w", encoding="utf-8"), indent=2)
print("DECK_DETAILED_OK", deck_path, len(narr), "slides + deck-narration.json")
