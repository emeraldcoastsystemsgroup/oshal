"""Generate a REAL graphical PowerPoint deck from the trading-day data (recap-data.json).
   Navy/electric OSHAL theme, big numbers, a winners bar chart. -> out/trading-recap.pptx
   Usage: python make-deck.py [recap-data.json]"""
import json, os, sys
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION

HERE = os.path.dirname(os.path.abspath(__file__))
OUT = os.path.join(HERE, "out")
data = json.load(open(sys.argv[1] if len(sys.argv) > 1 else os.path.join(OUT, "recap-data.json"), encoding="utf-8"))

NAVY = RGBColor(0x07, 0x0D, 0x1F); NAVY2 = RGBColor(0x0C, 0x17, 0x33)
CYAN = RGBColor(0x34, 0xE0, 0xFF); GREEN = RGBColor(0x27, 0xE0, 0x8A); RED = RGBColor(0xFF, 0x5D, 0x6C)
WHITE = RGBColor(0xEA, 0xF2, 0xFF); MUTED = RGBColor(0x8A, 0xA0, 0xC8)

prs = Presentation()
prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
SW, SH = prs.slide_width, prs.slide_height

def slide():
    s = prs.slides.add_slide(BLANK)
    s.background.fill.solid(); s.background.fill.fore_color.rgb = NAVY
    return s

def box(s, x, y, w, h, text, size, color, bold=True, align=PP_ALIGN.LEFT, font="Arial", anchor=MSO_ANCHOR.TOP):
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h)); tf = tb.text_frame
    tf.word_wrap = True; tf.vertical_anchor = anchor
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    f = r.font; f.size = Pt(size); f.bold = bold; f.color.rgb = color; f.name = font
    return tb

def bar(s, x, y, w, h, color=CYAN):
    sh = s.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb = color; sh.line.fill.background(); return sh

def brandbar(s):
    bar(s, 0, 0, 13.333, 0.12, CYAN)
    box(s, 0.6, 0.35, 6, 0.6, "oshal", 30, CYAN)
    box(s, 2.4, 0.5, 6, 0.4, "AUTONOMOUS TRADING DESK", 14, MUTED)
    box(s, 9.2, 0.5, 3.6, 0.4, data.get("date", ""), 14, MUTED, align=PP_ALIGN.RIGHT)

def money(n):
    s = "+" if n >= 0 else "-"; return f"{s}${abs(n):,.2f}"
def money0(n):
    s = "+" if n >= 0 else "-"; return f"{s}${abs(n):,.0f}"

# Slide 1 — title
s = slide()
box(s, 0, 2.4, 13.333, 1.6, "oshal", 110, CYAN, align=PP_ALIGN.CENTER)
box(s, 0, 4.0, 13.333, 0.8, "DAILY TRADING RECAP", 30, WHITE, align=PP_ALIGN.CENTER)
box(s, 0, 4.8, 13.333, 0.6, data.get("date", ""), 22, MUTED, align=PP_ALIGN.CENTER)

# Slide 2 — daily performance
s = slide(); brandbar(s)
box(s, 0.6, 1.4, 8, 0.5, "DAILY PERFORMANCE", 22, CYAN)
box(s, 0.5, 1.9, 12, 2.0, money(data["pl"]), 96, GREEN if data["pl"] >= 0 else RED)
box(s, 0.6, 3.9, 8, 0.6, f"Net day P/L  ·  {'+' if data['pct']>=0 else ''}{data['pct']:.2f}%", 24, MUTED)
box(s, 0.6, 5.2, 4, 0.4, "ACCOUNT EQUITY", 16, MUTED)
box(s, 0.6, 5.6, 4, 0.8, f"${data['equity']:,.0f}", 40, WHITE)
box(s, 5.0, 5.2, 4, 0.4, "OPEN UNREALIZED", 16, MUTED)
box(s, 5.0, 5.6, 4, 0.8, money0(data["unrealized"]), 40, GREEN if data["unrealized"] >= 0 else RED)
box(s, 9.4, 5.2, 3.4, 0.4, "MODE", 16, MUTED)
box(s, 9.4, 5.6, 3.4, 0.8, "Fully autonomous", 26, CYAN)

# Slide 3 — top winners (real bar chart)
s = slide(); brandbar(s)
box(s, 0.6, 1.4, 10, 0.5, "TOP OPEN WINNERS", 22, CYAN)
box(s, 0.6, 1.9, 11, 0.7, "Unrealized gains leading the book", 30, WHITE)
wins = data.get("winners", [])[:4]
if wins:
    cd = CategoryChartData(); cd.categories = [w[0] for w in wins]
    cd.add_series("Unrealized %", tuple(float(w[1]) for w in wins))
    gf = s.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(0.6), Inches(2.8), Inches(12), Inches(4.2), cd)
    ch = gf.chart; ch.has_legend = False; ch.has_title = False
    plot = ch.plots[0]; plot.gap_width = 60
    ser = plot.series[0]; ser.format.fill.solid(); ser.format.fill.fore_color.rgb = CYAN
    for ax in (ch.category_axis, ch.value_axis):
        ax.tick_labels.font.size = Pt(20); ax.tick_labels.font.color.rgb = WHITE
        ax.format.line.color.rgb = MUTED

# Slide 4 — the day
s = slide(); brandbar(s)
box(s, 0.6, 1.4, 10, 0.5, "THE DAY", 22, CYAN)
fills = data.get("fills", 0)
box(s, 0.6, 1.9, 12, 1.4, f"{fills} trades filled" if fills else "No new trades — holding the book", 56, WHITE)
box(s, 0.6, 3.5, 6, 0.4, "MORNING · LOCKED GAINS", 16, RED)
box(s, 0.6, 3.9, 6, 1.2, "  ".join(data.get("sells", [])) or "—", 30, WHITE)
box(s, 0.6, 5.1, 6, 0.4, "AFTERNOON · ROTATION", 16, GREEN)
box(s, 0.6, 5.5, 12, 1.2, "  ".join(data.get("buys", [])) or "—", 30, WHITE)

# Slide 5 — book at close
s = slide(); brandbar(s)
box(s, 0.6, 1.4, 10, 0.5, "BOOK AT CLOSE", 22, CYAN)
box(s, 0.6, 2.2, 6, 2.2, str(data.get("positions", 0)), 150, CYAN)
box(s, 0.6, 4.6, 6, 0.6, "OPEN POSITIONS", 22, MUTED)
box(s, 7.0, 2.4, 6, 0.5, "OPEN UNREALIZED", 18, MUTED)
box(s, 7.0, 2.8, 6, 1.0, money0(data["unrealized"]), 52, GREEN if data["unrealized"] >= 0 else RED)
box(s, 7.0, 4.2, 6, 0.5, "LEADERS", 18, MUTED)
box(s, 7.0, 4.6, 6, 0.8, (data.get("leaders") or "").replace(" · ", "   "), 36, WHITE)

outp = os.path.join(OUT, "trading-recap.pptx")
prs.save(outp)
print("DECK_OK", outp, len(prs.slides.__iter__.__self__._sldIdLst), "slides")
